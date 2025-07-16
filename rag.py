import os
import re
import json
from typing import List, Any, Dict, Tuple, Optional, Union
from dataclasses import dataclass, field
from dotenv import load_dotenv

import pdfplumber
import pandas as pd
from pptx import Presentation
from pptx.shapes.base import BaseShape

from langchain_community.vectorstores import Chroma
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain.schema.document import Document
from langchain_core.messages import BaseMessage

from langgraph.graph import StateGraph, END


load_dotenv()
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
assert OPENAI_API_KEY, "Please set OPENAI_API_KEY in your .env file."

@dataclass
class GraphState:
    docs: List[Union[str, Document]] = field(default_factory=list)
    vectorstore: Any = None
    question: str = ""
    retrieved_docs: List[Document] = field(default_factory=list)
    answer: str = ""

    selected_indexes: List[str] = field(default_factory=list)
    vector_candidates: List[Dict] = field(default_factory=list)
    reranked_parents: List[Dict] = field(default_factory=list)
    final_context: str = ""
    structured_answer: Dict = field(default_factory=dict)

def parse_pdf(path: str) -> List[Document]:
    docs = []
    with pdfplumber.open(path) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                # Split the text into chunks of 300 tokens with 50 tokens overlap
                chunks = split_text_into_chunks(text.strip(), chunk_size=300, overlap=50)
                for chunk_idx, chunk in enumerate(chunks):
                    if chunk.strip():
                        docs.append(Document(
                            page_content=chunk.strip(),
                            metadata={
                                "source_file": os.path.basename(path), 
                                "type": "pdf", 
                                "page": page_num,
                                "chunk": chunk_idx + 1,
                                "total_chunks_on_page": len(chunks)
                            }
                        ))
    return docs

def split_text_into_chunks(text: str, chunk_size: int = 300, overlap: int = 50) -> List[str]:
    """
    Split text into chunks of specified token size with overlap.
    Uses simple word-based tokenization (approximation).
    """
    words = text.split()
    chunks = []
    
    if len(words) <= chunk_size:
        return [text]
    
    start = 0
    while start < len(words):
        end = min(start + chunk_size, len(words))
        chunk_words = words[start:end]
        chunk_text = " ".join(chunk_words)
        chunks.append(chunk_text)
        
        # If we've reached the end, break
        if end >= len(words):
            break
            
        # Move start position considering overlap
        start = end - overlap
    
    return chunks

def route_indexes(question: str) -> List[str]:
    """
    Determine which indexes to search based on the question content.
    Returns list of file types to search: ['pdf', 'pptx', 'xls']
    """
    question_lower = question.lower()
    
    # Check for specific file mentions or keywords
    if any(keyword in question_lower for keyword in ['pdf', 'report', 'tsmc', '年報', '財報']):
        return ['pdf']
    elif any(keyword in question_lower for keyword in ['pptx', 'presentation', 'slide', 'digitimes', 'ai', '伺服器']):
        return ['pptx'] 
    elif any(keyword in question_lower for keyword in ['excel', 'xls', 'financial', 'fs', 'consolidated']):
        return ['xls']
    else:
        # Query all indexes if no specific mention
        return ['pdf', 'pptx', 'xls']

def get_parent_id(doc: Document) -> str:
    """
    Generate a unique parent ID for grouping chunks from the same source unit.
    For chunked documents, we group by the original source unit (page/slide/sheet_rows).
    """
    metadata = doc.metadata
    if metadata['type'] == 'pdf':
        return f"pdf_page_{metadata['page']}"
    elif metadata['type'] == 'pptx':
        return f"pptx_slide_{metadata['slide']}"
    elif metadata['type'] == 'xls':
        # For XLS, we group by sheet and row range to maintain granularity
        return f"xls_{metadata['sheet']}_rows_{metadata['rows']}"
    return "unknown"

def aggregate_parents(candidates: List[Dict]) -> List[Dict]:
    """
    Collapse duplicate parents, keeping the highest vector score for each parent.
    """
    parent_groups = {}
    
    for candidate in candidates:
        parent_id = candidate['parent_id']
        if parent_id not in parent_groups:
            parent_groups[parent_id] = candidate
        else:
            # Keep the one with higher vector score
            if candidate['vector_score'] > parent_groups[parent_id]['vector_score']:
                parent_groups[parent_id] = candidate
    
    return list(parent_groups.values())

def llm_rerank_batch(texts: List[str], question: str) -> List[float]:
    """
    Send a batch of texts to GPT-4o-mini for relevance scoring.
    Returns list of scores between 0-1.
    """
    llm = ChatOpenAI(model="gpt-4o-mini", temperature=0)
    
    prompt = f"""You are a relevance scorer. Given a question and up to 3 text passages, rate each passage's relevance to answering the question.

Question: {question}

Passages:
"""
    for i, text in enumerate(texts, 1):
        prompt += f"\nPassage {i}:\n{text[:500]}...\n"
    
    prompt += """
Rate each passage from 0.0 to 1.0 (where 1.0 = highly relevant, 0.0 = not relevant).
Respond with ONLY a JSON array of scores, e.g., [0.8, 0.2, 0.9]
"""
    
    try:
        response = llm.invoke(prompt)
        response_content = response.content if isinstance(response, BaseMessage) else str(response)
        if isinstance(response_content, str):
            scores = json.loads(response_content)
            # Ensure we have the right number of scores
            while len(scores) < len(texts):
                scores.append(0.0)
            return scores[:len(texts)]
        else:
            return [0.5] * len(texts)
    except:
        # Fallback to neutral scores if parsing fails
        return [0.5] * len(texts)

def fuse_scores(vector_score: float, llm_score: float) -> float:
    """
    Combine vector and LLM scores with 0.3 weight for vector, 0.7 for LLM.
    """
    return 0.3 * vector_score + 0.7 * llm_score

def assemble_context(parents: List[Dict]) -> str:
    """
    Create final context string with parent IDs as prefixes.
    """
    context_parts = []
    for parent in parents:
        prefix = f"[{parent['parent_id']}]"
        content = parent['content']
        context_parts.append(f"{prefix}\n{content}")
    
    return "\n\n".join(context_parts)

def is_complex_query(question: str) -> bool:
    """
    Determine if a query requires complex query splitting.
    """
    complexity_indicators = [
        'compare', 'versus', 'vs', 'difference', 'contrast',
        'and', 'both', 'all', 'multiple', 'several',
        '比較', '對比', '差異', '和', '都', '多個'
    ]
    return any(indicator in question.lower() for indicator in complexity_indicators)

def split_complex_query(question: str) -> List[str]:
    """
    Use GPT-4o-mini to break down complex queries into sub-queries.
    """
    llm = ChatOpenAI(model="gpt-4o-mini", temperature=0)
    
    prompt = f"""Break down this complex question into 2-4 simpler sub-questions that can be answered independently.

Original question: {question}

Return ONLY a JSON array of sub-questions, e.g.:
["What is X?", "What is Y?", "How do X and Y compare?"]
"""
    
    try:
        response = llm.invoke(prompt)
        response_content = response.content if isinstance(response, BaseMessage) else str(response)
        if isinstance(response_content, str):
            sub_queries = json.loads(response_content)
            return sub_queries if isinstance(sub_queries, list) else [question]
        else:
            return [question]
    except:
        return [question]

def create_answer_prompt(context: str, question: str) -> str:
    """
    Create the structured prompt for answering with JSON schema.
    """
    system_message = """You are a knowledgeable assistant that provides thorough, accurate answers based on the given context. Follow these rules:

1. Base your answer ONLY on the provided context
2. Use step-by-step reasoning in your analysis
3. Cite specific sources using the provided IDs
4. Be precise and comprehensive
5. If information is insufficient, state so clearly
"""
    
    user_prompt = f"""CONTEXT:
{context}

QUESTION: {question}

Provide your response in the following JSON format:
{{
    "step_by_step_analysis": "Your detailed reasoning process here",
    "relevant_sources": ["list", "of", "source", "IDs", "used"],
    "final_answer": "Your comprehensive final answer here",
    "confidence_level": "high/medium/low based on available information"
}}
"""
    
    return system_message + "\n\n" + user_prompt

def validate_and_fix_json_response(response_text: str, question: str) -> Dict:
    """
    Parse JSON response and fix if needed.
    """
    try:
        return json.loads(response_text)
    except json.JSONDecodeError:
        # Try to extract JSON from the response
        json_match = re.search(r'\{.*\}', response_text, re.DOTALL)
        if json_match:
            try:
                return json.loads(json_match.group())
            except:
                pass
        
        # Fallback: create structured response from unstructured text
        return {
            "step_by_step_analysis": "Unable to parse structured analysis from response",
            "relevant_sources": [],
            "final_answer": response_text,
            "confidence_level": "low"
        }

def parse_pptx(path: str) -> List[Document]:
    docs = []
    prs = Presentation(path)
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            # Check if shape has text_frame and text
            try:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text_content = shape.text_frame.text
                    if text_content and text_content.strip():
                        slide_text.append(text_content)
            except:
                pass
            
            # Check if shape is a table
            try:
                if hasattr(shape, "table") and shape.table:
                    for row in shape.table.rows:
                        row_text = " | ".join([cell.text for cell in row.cells])
                        if row_text.strip():
                            slide_text.append(row_text)
            except:
                pass  # Skip if table parsing fails
        
        text = "\n".join(slide_text)
        if text.strip():
            # Split the slide text into chunks of 300 tokens with 50 tokens overlap
            chunks = split_text_into_chunks(text.strip(), chunk_size=300, overlap=50)
            for chunk_idx, chunk in enumerate(chunks):
                if chunk.strip():
                    docs.append(Document(
                        page_content=chunk.strip(),
                        metadata={
                            "source_file": os.path.basename(path), 
                            "type": "pptx", 
                            "slide": slide_num,
                            "chunk": chunk_idx + 1,
                            "total_chunks_on_slide": len(chunks)
                        }
                    ))
    return docs

def parse_xls(path: str) -> List[Document]:
    docs = []
    try:
        xls = pd.ExcelFile(path)
        for sheet in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet).astype(str)
            
            if len(df) == 0:
                continue
                
            # Get column headers for better context
            headers = ", ".join(df.columns.tolist())
            
            # Create chunks that preserve row integrity
            chunks = split_excel_into_chunks(df, headers, chunk_size=300, overlap=50)
            
            for chunk_idx, (chunk_text, start_row, end_row) in enumerate(chunks):
                if chunk_text.strip():
                    docs.append(Document(
                        page_content=chunk_text.strip(),
                        metadata={
                            "source_file": os.path.basename(path), 
                            "type": "xls", 
                            "sheet": sheet, 
                            "chunk": chunk_idx + 1,
                            "total_chunks_on_sheet": len(chunks),
                            "rows": f"{start_row}-{end_row}",
                            "headers": headers
                        }
                    ))
    except Exception as e:
        print(f"Error parsing Excel file {path}: {e}")
    return docs

def split_excel_into_chunks(df: pd.DataFrame, headers: str, chunk_size: int = 300, overlap: int = 50) -> List[Tuple[str, int, int]]:
    """
    Split Excel dataframe into chunks while preserving row integrity.
    Returns list of (chunk_text, start_row, end_row) tuples.
    """
    chunks = []
    total_rows = len(df)
    
    if total_rows == 0:
        return chunks
    
    # Estimate tokens per row (rough approximation)
    sample_rows = min(5, total_rows)
    sample_text = "\n".join([", ".join(row) for row in df.head(sample_rows).values])
    avg_tokens_per_row = max(1, len(sample_text.split()) / sample_rows)
    
    # Calculate rows per chunk
    rows_per_chunk = max(1, int(chunk_size / avg_tokens_per_row))
    overlap_rows = max(0, int(overlap / avg_tokens_per_row))
    
    start_row = 0
    chunk_idx = 0
    
    while start_row < total_rows:
        end_row = min(start_row + rows_per_chunk, total_rows)
        
        # Extract chunk data
        chunk_df = df.iloc[start_row:end_row]
        
        # Create chunk text with headers
        chunk_text_parts = [f"Headers: {headers}"]
        chunk_text_parts.extend([", ".join(row) for row in chunk_df.values])
        chunk_text = "\n".join(chunk_text_parts)
        
        chunks.append((chunk_text, start_row + 1, end_row))  # 1-indexed for user readability
        
        # Move to next chunk with overlap
        if end_row >= total_rows:
            break
        start_row = end_row - overlap_rows
        
        # Ensure we don't go backwards
        if start_row <= chunk_idx * rows_per_chunk:
            start_row = end_row
    
    return chunks

def ingest_node(state: GraphState) -> GraphState:
    # Only parse documents if docs field is a list of file paths (not Document objects)
    if state.docs and isinstance(state.docs[0], str):
        pdf_docs = parse_pdf(str(state.docs[0])) if len(state.docs) > 0 else []
        pptx_docs = parse_pptx(str(state.docs[1])) if len(state.docs) > 1 else []
        xls_docs = parse_xls(str(state.docs[2])) if len(state.docs) > 2 else []
        docs = pdf_docs + pptx_docs + xls_docs
    else:
        docs = [doc for doc in state.docs if isinstance(doc, Document)]
    
    docs_final: List[Union[str, Document]] = docs  # Explicit type annotation
    return GraphState(
        docs=docs_final,
        vectorstore=state.vectorstore,
        question=state.question
    )

def embed_node(state: GraphState) -> GraphState:
    # Only build the vectorstore if it doesn't exist yet
    if state.vectorstore is None:
        embeddings = OpenAIEmbeddings(model="text-embedding-3-large")
        # Ensure we only pass Document objects to Chroma
        document_list = [doc for doc in state.docs if isinstance(doc, Document)]
        vectorstore = Chroma.from_documents(
            document_list,
            embedding=embeddings,
            persist_directory="chroma_db_new"
        )
    else:
        vectorstore = state.vectorstore
    return GraphState(
        docs=state.docs,
        vectorstore=vectorstore,
        question=state.question
    )

def retrieval_node(state: GraphState) -> GraphState:
    """
    Enhanced retrieval with index routing, parent aggregation, and LLM reranking.
    """
    # Step 1: Index routing
    selected_indexes = route_indexes(state.question)
    
    # Step 2: Vector search - get top 30 from each selected index
    all_candidates = []
    docs_with_scores = state.vectorstore.similarity_search_with_score(state.question, k=30)
    
    for doc, score in docs_with_scores:
        if doc.metadata['type'] in selected_indexes:
            parent_id = get_parent_id(doc)
            candidate = {
                'doc': doc,
                'parent_id': parent_id,
                'vector_score': 1 - score,  # Convert distance to similarity
                'content': doc.page_content,
                'metadata': doc.metadata
            }
            all_candidates.append(candidate)
    
    # Step 3: Parent aggregation
    aggregated_parents = aggregate_parents(all_candidates)
    
    # Step 4: LLM reranking in batches of 3
    reranked_parents = []
    for i in range(0, len(aggregated_parents), 3):
        batch = aggregated_parents[i:i+3]
        texts = [parent['content'] for parent in batch]
        llm_scores = llm_rerank_batch(texts, state.question)
        
        for j, parent in enumerate(batch):
            # Step 5: Score fusion
            fused_score = fuse_scores(parent['vector_score'], llm_scores[j])
            parent['fused_score'] = fused_score
            parent['llm_score'] = llm_scores[j]
            reranked_parents.append(parent)
    
    # Keep top 10 highest-scoring parents
    reranked_parents.sort(key=lambda x: x['fused_score'], reverse=True)
    top_parents = reranked_parents[:10]
    
    # Step 6: Context assembly
    final_context = assemble_context(top_parents)
    
    # Convert back to Documents for compatibility
    retrieved_docs = [parent['doc'] for parent in top_parents]
    
    return GraphState(
        docs=state.docs,
        vectorstore=state.vectorstore,
        question=state.question,
        retrieved_docs=retrieved_docs,
        selected_indexes=selected_indexes,
        vector_candidates=all_candidates,
        reranked_parents=top_parents,
        final_context=final_context
    )

def rag_node(state: GraphState) -> GraphState:
    """
    Enhanced answering with structured output, chain-of-thought, and complex query handling.
    """
    # Check if this is a complex query that needs splitting
    if is_complex_query(state.question):
        return handle_complex_query(state)
    else:
        return handle_simple_query(state)

def handle_simple_query(state: GraphState) -> GraphState:
    """
    Handle simple queries with structured JSON output.
    """
    llm = ChatOpenAI(model="gpt-4o", temperature=0.3)
    
    # Use the assembled context from retrieval
    context = state.final_context
    prompt = create_answer_prompt(context, state.question)
    
    # Get structured response
    response = llm.invoke(prompt)
    response_content = response.content if isinstance(response, BaseMessage) else str(response)
    response_str = str(response_content) if not isinstance(response_content, str) else response_content
    structured_answer = validate_and_fix_json_response(response_str, state.question)
    
    # Extract final answer for backward compatibility
    final_answer = structured_answer.get('final_answer', response_str)
    
    return GraphState(
        docs=state.docs,
        vectorstore=state.vectorstore,
        question=state.question,
        retrieved_docs=state.retrieved_docs,
        answer=final_answer,
        selected_indexes=state.selected_indexes,
        vector_candidates=state.vector_candidates,
        reranked_parents=state.reranked_parents,
        final_context=state.final_context,
        structured_answer=structured_answer
    )

def handle_complex_query(state: GraphState) -> GraphState:
    """
    Handle complex queries by splitting into sub-queries and synthesizing results.
    """
    # Split the complex query
    sub_queries = split_complex_query(state.question)
    
    # Process each sub-query
    sub_results = []
    for sub_query in sub_queries:
        # Create a new state for the sub-query
        sub_state = GraphState(
            docs=state.docs,
            vectorstore=state.vectorstore,
            question=sub_query
        )
        
        # Run retrieval for the sub-query
        sub_state_after_retrieval = retrieval_node(sub_state)
        
        # Get answer for the sub-query
        sub_state_final = handle_simple_query(sub_state_after_retrieval)
        
        sub_results.append({
            'question': sub_query,
            'answer': sub_state_final.structured_answer.get('final_answer', ''),
            'sources': sub_state_final.structured_answer.get('relevant_sources', [])
        })
    
    # Synthesize all sub-results
    synthesis_answer = synthesize_sub_results(state.question, sub_results)
    
    return GraphState(
        docs=state.docs,
        vectorstore=state.vectorstore,
        question=state.question,
        retrieved_docs=state.retrieved_docs,
        answer=synthesis_answer.get('final_answer', ''),
        selected_indexes=state.selected_indexes,
        vector_candidates=state.vector_candidates,
        reranked_parents=state.reranked_parents,
        final_context=state.final_context,
        structured_answer=synthesis_answer
    )

def synthesize_sub_results(original_question: str, sub_results: List[Dict]) -> Dict:
    """
    Synthesize multiple sub-query results into a comprehensive answer.
    """
    llm = ChatOpenAI(model="gpt-4o", temperature=0.3)
    
    sub_results_text = "\n\n".join([
        f"Sub-question: {result['question']}\nAnswer: {result['answer']}\nSources: {result['sources']}"
        for result in sub_results
    ])
    
    prompt = f"""Based on the following sub-query results, provide a comprehensive answer to the original question.

Original Question: {original_question}

Sub-query Results:
{sub_results_text}

Provide your response in the following JSON format:
{{
    "step_by_step_analysis": "How you synthesized the sub-results",
    "relevant_sources": ["combined", "list", "of", "all", "sources"],
    "final_answer": "Comprehensive answer combining all sub-results",
    "confidence_level": "high/medium/low"
}}
"""
    
    response = llm.invoke(prompt)
    response_content = response.content if isinstance(response, BaseMessage) else str(response)
    response_str = str(response_content) if not isinstance(response_content, str) else response_content
    return validate_and_fix_json_response(response_str, original_question)

def log_node(state: GraphState) -> GraphState:
    """
    Enhanced logging with detailed retrieval and answering metrics.
    """
    row = {
        "question": state.question,
        "answer": state.answer,
        "selected_indexes": str(state.selected_indexes),
        "num_candidates": len(state.vector_candidates),
        "num_final_parents": len(state.reranked_parents),
        "confidence_level": state.structured_answer.get('confidence_level', 'unknown'),
        "relevant_sources": str(state.structured_answer.get('relevant_sources', [])),
        "step_by_step_analysis": state.structured_answer.get('step_by_step_analysis', '')[:200] + "..." if len(state.structured_answer.get('step_by_step_analysis', '')) > 200 else state.structured_answer.get('step_by_step_analysis', ''),
        "retrieved_docs_metadata": str([doc.metadata for doc in state.retrieved_docs])
    }
    df = pd.DataFrame([row])
    if not os.path.exists("rag_qa_log.csv"):
        df.to_csv("rag_qa_log.csv", index=False)
    else:
        df.to_csv("rag_qa_log.csv", mode="a", header=False, index=False)
    return state


init_workflow = StateGraph(GraphState)
init_workflow.add_node("ingest", ingest_node)
init_workflow.add_node("embed", embed_node)
init_workflow.set_entry_point("ingest")
init_workflow.add_edge("ingest", "embed")
init_workflow.add_edge("embed", END)
init_graph = init_workflow.compile()


query_workflow = StateGraph(GraphState)
query_workflow.add_node("retrieval", retrieval_node)
query_workflow.add_node("rag", rag_node)
query_workflow.add_node("log", log_node)
query_workflow.set_entry_point("retrieval")
query_workflow.add_edge("retrieval", "rag")
query_workflow.add_edge("rag", "log")
query_workflow.add_edge("log", END)
query_graph = query_workflow.compile()

if __name__ == "__main__":
    # Specify your file paths
    PDF_PATH = "./tsmc_2024_yearly_report.pdf"
    PPTX_PATH = "./無備忘錄版_digitimes_2025年全球AI伺服器出貨將達181萬台　高階機種採購不再集中於四大CSP.pptx"
    XLS_PATH = "./excel_FS-Consolidated_1Q25.xls"
    print("RAG system is going to be ready. Enter your questions one by one. Type 'quit' to exit....\n")

    # Initial build: parse and embed only once
    initial_state = GraphState(
        docs=[PDF_PATH, PPTX_PATH, XLS_PATH],
        question=""
    )
    state_after_embed =init_graph.invoke(initial_state)
    # After this, reuse the parsed docs and vectorstore for all following queries
    print("Documents parsed and embedded successfully. You can now ask questions.\n")
    docs = state_after_embed["docs"]
    vectorstore = state_after_embed["vectorstore"]

    while True:
        question = input("\n" + "=" * 20 + " Question (or 'quit' to exit) " + "=" * 20 + "\n")
        
        if question.strip().lower() == "quit":
            print("Exiting RAG system....")
            break
        state = GraphState(
            docs=docs,
            vectorstore=vectorstore,
            question=question
        )
        result = query_graph.invoke(state)
        
        print("\n" + "=" * 5 + " RAG RESULTS " + "=" * 5)
        print(f"Selected Indexes: {result['selected_indexes']}")
        print(f"Confidence Level: {result['structured_answer'].get('confidence_level', 'unknown')}")
        print(f"\n--- Step-by-Step Analysis ---")
        print(result['structured_answer'].get('step_by_step_analysis', 'No analysis available'))
        print(f"\n--- Relevant Sources ---")
        print(result['structured_answer'].get('relevant_sources', []))
        print(f"\n" + "=" * 20 + "Final Answer" + "=" * 20)
        print(result["answer"])
